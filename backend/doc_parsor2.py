import os
import json
import logging
from dataclasses import dataclass, asdict
from datetime import datetime
from typing import List, Dict, Any, Tuple, Optional
import torch
import numpy as np
import cv2
from PIL import Image
import io
import pymupdf as fitz
from transformers import VisionEncoderDecoderModel, ViTImageProcessor, AutoTokenizer
from llama_index.core import Document, Settings, StorageContext, VectorStoreIndex
from llama_index.embeddings.openai import OpenAIEmbedding
from llama_index.vector_stores.pinecone import PineconeVectorStore
from llama_index.llms.openai import OpenAI
from llama_parse import LlamaParse
from unstructured.partition.pdf import partition_pdf
from unstructured.partition.auto import partition

# Set up logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('document_processing.log'),
        logging.StreamHandler()
    ]
)

logger = logging.getLogger(__name__)

@dataclass
class ProcessingMetrics:
    start_time: datetime
    end_time: Optional[datetime] = None
    document_count: int = 0
    total_tokens: int = 0
    total_characters: int = 0
    processed_files: List[str] = None
    errors: List[Dict] = None
    
    def to_dict(self):
        return asdict(self)
    
    def log_error(self, filename: str, error: Exception):
        if self.errors is None:
            self.errors = []
        self.errors.append({
            "filename": filename,
            "error": str(error),
            "timestamp": datetime.now().isoformat()
        })

@dataclass
class DocumentElement:
    content: Any
    element_type: str
    page_num: int
    bbox: Dict
    source_file: str
    element_id: str
    citations: List[str] = None
    
    def to_document(self) -> Document:
        metadata = {
            "type": self.element_type,
            "page_num": self.page_num,
            "bbox": self.bbox,
            "source": self.source_file,
            "element_id": self.element_id,
            "citations": self.citations
        }
        
        if self.element_type == "table":
            metadata["table_path"] = self.content["table_path"]
            text = f"Table with caption: {self.content['caption']}\nColumns: {self.content['columns']}"
        elif self.element_type == "image":
            metadata["image_path"] = self.content["image_path"]
            text = f"Image with caption: {self.content['caption']}"
        else:
            text = self.content
            
        return Document(text=text, metadata=metadata, id_=self.element_id)

class DocumentProcessor:
    def __init__(self):
        self.metrics = ProcessingMetrics(
            start_time=datetime.now(),
            processed_files=[],
            errors=[]
        )
        
        # Initialize models
        self.vision_model = VisionEncoderDecoderModel.from_pretrained("nlpconnect/vit-gpt2-image-captioning")
        self.feature_extractor = ViTImageProcessor.from_pretrained("nlpconnect/vit-gpt2-image-captioning")
        self.tokenizer = AutoTokenizer.from_pretrained("nlpconnect/vit-gpt2-image-captioning")
        
        self.device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
        self.vision_model.to(self.device)
        
        logger.info(f"Initialized DocumentProcessor with device: {self.device}")

    def process_text_block(self, text_block: tuple, page_num: int, source_file: str, block_id: int) -> DocumentElement:
        """Process a text block from PDF."""
        logger.debug(f"Processing text block {block_id} from page {page_num}")
        
        bbox = {
            "x1": text_block[0], 
            "y1": text_block[1],
            "x2": text_block[2], 
            "y2": text_block[3]
        }
        
        self.metrics.total_characters += len(text_block[4])
        # Approximate token count
        self.metrics.total_tokens += len(text_block[4].split())
        
        return DocumentElement(
            content=text_block[4],
            element_type="text",
            page_num=page_num,
            bbox=bbox,
            source_file=source_file,
            element_id=f"{source_file}-page{page_num}-block{block_id}"
        )

    def process_table(self, table, page_num: int, source_file: str, table_id: int, 
                     save_dir: str) -> DocumentElement:
        """Process a table from PDF."""
        logger.debug(f"Processing table {table_id} from page {page_num}")
        
        df = table.to_pandas()
        table_path = os.path.join(save_dir, f"table-{table_id}-page-{page_num}.xlsx")
        df.to_excel(table_path, index=False)
        
        content = {
            "table_path": table_path,
            "caption": f"Table extracted from page {page_num}",
            "columns": ", ".join(df.columns.tolist())
        }
        
        # Update metrics
        self.metrics.total_characters += sum(df.astype(str).sum().sum())
        self.metrics.total_tokens += len(df.columns) * df.shape[0]  # Approximate
        
        return DocumentElement(
            content=content,
            element_type="table",
            page_num=page_num,
            bbox=table.bbox,
            source_file=source_file,
            element_id=f"{source_file}-page{page_num}-table{table_id}"
        )

    def describe_image(self, image_content: bytes) -> str:
        """Generate a text description of the image using a pre-trained model."""
        logger.debug("Generating image description")
        
        image = Image.open(io.BytesIO(image_content)).convert('RGB')
        pixel_values = self.feature_extractor(images=[image], return_tensors="pt").pixel_values
        pixel_values = pixel_values.to(self.device)

        output_ids = self.vision_model.generate(
            pixel_values,
            max_length=16,
            num_beams=4
        )

        preds = self.tokenizer.batch_decode(output_ids, skip_special_tokens=True)
        description = preds[0].strip()
        
        # Update metrics
        self.metrics.total_characters += len(description)
        self.metrics.total_tokens += len(description.split())
        
        return description

    def is_graph(self, image_data: bytes) -> bool:
        """Detect if the image contains a graph."""
        logger.debug("Checking if image is a graph")
        
        nparr = np.frombuffer(image_data, np.uint8)
        img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        edges = cv2.Canny(gray, 50, 150, apertureSize=3)
        lines = cv2.HoughLines(edges, 1, np.pi/180, threshold=100)
        
        return lines is not None and len(lines) > 10

    def process_graph(self, image_data: bytes) -> str:
        """Process a graph image to extract basic information."""
        logger.debug("Processing graph")
        
        nparr = np.frombuffer(image_data, np.uint8)
        img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        edges = cv2.Canny(gray, 50, 150, apertureSize=3)
        lines = cv2.HoughLines(edges, 1, np.pi/180, threshold=100)
        
        horizontal_lines = 0
        vertical_lines = 0
        
        if lines is not None:
            for line in lines:
                rho, theta = line[0]
                if np.abs(theta) < np.pi / 4 or np.abs(theta - np.pi) < np.pi / 4:
                    vertical_lines += 1
                else:
                    horizontal_lines += 1
        
        graph_type = "bar graph" if vertical_lines > horizontal_lines else \
                    "line graph" if horizontal_lines > vertical_lines else \
                    "scatter plot"
        
        description = f"This image appears to be a {graph_type} with approximately {vertical_lines} vertical lines and {horizontal_lines} horizontal lines."
        
        # Update metrics
        self.metrics.total_characters += len(description)
        self.metrics.total_tokens += len(description.split())
        
        return description

    def parse_all_tables(self, filename: str, page, pagenum: int, text_blocks, ongoing_tables) -> Tuple[List[Document], List, Any]:
        """Extract tables from a PDF page."""
        logger.debug(f"Parsing tables from page {pagenum}")
        
        table_docs = []
        table_bboxes = []
        
        try:
            tables = page.find_tables(
                horizontal_strategy="lines_strict",
                vertical_strategy="lines_strict"
            )
            
            tablerefdir = os.path.join(os.getcwd(), "vectorstore/table_references")
            os.makedirs(tablerefdir, exist_ok=True)
            
            for tab in tables:
                if not tab.header.external:
                    continue
                
                pandas_df = tab.to_pandas()
                df_xlsx_path = os.path.join(
                    tablerefdir,
                    f"table{len(table_docs)+1}-page{pagenum}.xlsx"
                )
                pandas_df.to_excel(df_xlsx_path)
                
                bbox = fitz.Rect(tab.bbox)
                table_bboxes.append(bbox)
                
                # Extract surrounding text
                before_text, after_text = self.extract_text_around_item(
                    text_blocks,
                    bbox,
                    page.rect.height
                )
                
                # Save table image
                table_img = page.get_pixmap(clip=bbox)
                table_img_path = os.path.join(
                    tablerefdir,
                    f"table{len(table_docs)+1}-page{pagenum}.jpg"
                )
                table_img.save(table_img_path)
                
                description = self.process_graph(table_img.tobytes())
                caption = before_text.replace("\n", " ") + description + after_text.replace("\n", " ")
                if before_text == "" and after_text == "":
                    caption = " ".join(tab.header.names)
                
                table_metadata = {
                    "source": f"{filename[:-4]}-page{pagenum}-table{len(table_docs)+1}",
                    "excel_path": df_xlsx_path,
                    "image_path": table_img_path,
                    "caption": caption,
                    "type": "table",
                    "page_num": pagenum,
                    "columns": list(pandas_df.columns.values),
                    "markdown_table": pandas_df.to_markdown(),
                    "html_table": pandas_df.to_html()
                }
                
                doc_text = (
                    f"Table Reference: {table_metadata['markdown_table']}\n"
                    f"Caption: {caption}\n"
                    f"Excel File: {df_xlsx_path}\n"
                    f"Table Image: ![Table]({table_img_path})"
                )
                
                table_docs.append(Document(
                    text=doc_text,
                    metadata=table_metadata
                ))
                
                # Update metrics
                self.metrics.total_characters += len(doc_text)
                self.metrics.total_tokens += len(doc_text.split())
                
        except Exception as e:
            logger.error(f"Error processing tables on page {pagenum}: {e}")
            self.metrics.log_error(filename, e)
        
        return table_docs, table_bboxes, ongoing_tables

    def load_data_from_directory(self, directory: str) -> List[Document]:
        """Load and process multiple file types from a directory."""
        logger.info(f"Processing directory: {directory}")
        
        documents = []
        
        for filename in os.listdir(directory):
            filepath = os.path.join(directory, filename)
            file_extension = os.path.splitext(filename.lower())[1]
            
            logger.info(f"Processing file: {filename}")
            self.metrics.processed_files.append(filename)
            
            try:
                if file_extension in ('.png', '.jpg', '.jpeg'):
                    with open(filepath, "rb") as image_file:
                        image_content = image_file.read()
                        image_text = self.describe_image(image_content)
                        doc = Document(
                            text=image_text,
                            metadata={"source": filename, "type": "image"}
                        )
                        documents.append(doc)
                
                elif file_extension == '.pdf':
                    with open(filepath, "rb") as pdf_file:
                        pdf_documents = self.get_pdf_documents(pdf_file)
                        documents.extend(pdf_documents)
                
                elif file_extension in ('.ppt', '.pptx'):
                    ppt_documents = self.process_ppt_file(filepath)
                    documents.extend(ppt_documents)
                
                else:
                    with open(filepath, "r", encoding="utf-8") as text_file:
                        text = text_file.read()
                        doc = Document(
                            text=text,
                            metadata={"source": filename, "type": "text"}
                        )
                        documents.append(doc)
                
                self.metrics.document_count += 1
                
            except Exception as e:
                logger.error(f"Error processing {filename}: {e}")
                self.metrics.log_error(filename, e)
                continue
        
        return documents

    def export_metrics(self, output_path: str = "processing_metrics.json"):
        """Export processing metrics to JSON file."""
        self.metrics.end_time = datetime.now()
        
        with open(output_path, 'w') as f:
            json.dump(self.metrics.to_dict(), f, indent=2, default=str)
        
        logger.info(f"Metrics exported to {output_path}")
        
        return self.metrics.to_dict()

    # Helper method that needs to be implemented
    # ... [Previous code remains the same until extract_text_around_item] ...

    def extract_text_around_item(self, text_blocks, bbox, page_height, context_distance: float = 50.0) -> Tuple[str, str]:
        """
        Extract text before and after an item on the page within a specified distance.
        
        Args:
            text_blocks: List of text blocks from the page
            bbox: Bounding box of the item (table/image)
            page_height: Height of the page
            context_distance: Distance in points to look for context (default: 50.0)
            
        Returns:
            Tuple[str, str]: (text before item, text after item)
        """
        logger.debug("Extracting text around item")
        
        # Convert bbox to coordinates if it's not already
        if isinstance(bbox, dict):
            item_top = bbox['y1']
            item_bottom = bbox['y2']
        else:
            item_top = bbox.y0
            item_bottom = bbox.y1
            
        before_text = []
        after_text = []
        
        for block in text_blocks:
            # Get the vertical position of the text block
            block_bottom = block[3]  # y2
            block_top = block[1]     # y1
            block_text = block[4]    # text content
            
            # Check if text is before the item
            if block_bottom < item_top and (item_top - block_bottom) <= context_distance:
                before_text.append(block_text)
                
            # Check if text is after the item
            elif block_top > item_bottom and (block_top - item_bottom) <= context_distance:
                after_text.append(block_text)
        
        return " ".join(before_text), " ".join(after_text)

    def get_pdf_documents(self, pdf_file) -> List[Document]:
        """
        Process PDF file and return list of documents.
        
        Args:
            pdf_file: File object of the PDF
            
        Returns:
            List[Document]: List of processed documents
        """
        logger.info("Processing PDF file")
        
        try:
            # Load PDF
            pdf_document = fitz.open(stream=pdf_file.read())
            documents = []
            
            for page_num in range(len(pdf_document)):
                page = pdf_document[page_num]
                
                # Get text blocks
                text_blocks = page.get_text("blocks")
                
                # Process tables first
                table_docs, table_bboxes, _ = self.parse_all_tables(
                    pdf_file.name,
                    page,
                    page_num + 1,
                    text_blocks,
                    None
                )
                documents.extend(table_docs)
                
                # Process remaining text blocks
                for block_id, block in enumerate(text_blocks):
                    # Skip if block overlaps with any table
                    block_rect = fitz.Rect(block[:4])
                    skip_block = any(block_rect.intersects(table_bbox) for table_bbox in table_bboxes)
                    
                    if not skip_block:
                        doc_element = self.process_text_block(
                            block,
                            page_num + 1,
                            pdf_file.name,
                            block_id
                        )
                        documents.append(doc_element.to_document())
                
                # Process images
                images = page.get_images(full=True)
                for img_id, img_info in enumerate(images):
                    xref = img_info[0]
                    base_image = pdf_document.extract_image(xref)
                    
                    if base_image:
                        image_data = base_image["image"]
                        
                        # Process as graph if detected
                        if self.is_graph(image_data):
                            description = self.process_graph(image_data)
                        else:
                            description = self.describe_image(image_data)
                        
                        # Create image document
                        doc = Document(
                            text=description,
                            metadata={
                                "source": pdf_file.name,
                                "type": "image",
                                "page_num": page_num + 1,
                                "image_id": img_id
                            }
                        )
                        documents.append(doc)
            
            return documents
            
        except Exception as e:
            logger.error(f"Error processing PDF: {e}")
            self.metrics.log_error(pdf_file.name, e)
            return []

    def process_ppt_file(self, filepath: str) -> List[Document]:
        """
        Process PowerPoint file and return list of documents.
        
        Args:
            filepath: Path to the PowerPoint file
            
        Returns:
            List[Document]: List of processed documents
        """
        logger.info(f"Processing PowerPoint file: {filepath}")
        
        try:
            from pptx import Presentation
            
            documents = []
            prs = Presentation(filepath)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                # Process text content
                text_content = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content.append(shape.text.strip())
                        
                        # Update metrics
                        self.metrics.total_characters += len(shape.text)
                        self.metrics.total_tokens += len(shape.text.split())
                
                if text_content:
                    doc = Document(
                        text="\n".join(text_content),
                        metadata={
                            "source": filepath,
                            "type": "presentation",
                            "slide_num": slide_num
                        }
                    )
                    documents.append(doc)
                
                # Process images in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "image"):
                        try:
                            # Save image temporarily
                            image_bytes = shape.image.blob
                            
                            # Process as graph if detected
                            if self.is_graph(image_bytes):
                                description = self.process_graph(image_bytes)
                            else:
                                description = self.describe_image(image_bytes)
                            
                            doc = Document(
                                text=description,
                                metadata={
                                    "source": filepath,
                                    "type": "presentation_image",
                                    "slide_num": slide_num
                                }
                            )
                            documents.append(doc)
                            
                        except Exception as e:
                            logger.warning(f"Error processing image in slide {slide_num}: {e}")
                            continue
            
            return documents
            
        except Exception as e:
            logger.error(f"Error processing PowerPoint file: {e}")
            self.metrics.log_error(filepath, e)
            return []

if __name__ == "__main__":
    # Example usage
    processor = DocumentProcessor()
    
    # Process a directory of documents
    input_directory = "path/to/your/documents"
    documents = processor.load_data_from_directory(input_directory)
    
    # Export processing metrics
    metrics = processor.export_metrics()
    
    # Print summary
    print(f"Processed {metrics['document_count']} documents")
    print(f"Total characters: {metrics['total_characters']}")
    print(f"Total tokens: {metrics['total_tokens']}")
    print(f"Processing time: {metrics['end_time'] - metrics['start_time']}")
    
    if metrics['errors']:
        print("\nErrors encountered:")
        for error in metrics['errors']:
            print(f"- {error['filename']}: {error['error']}")